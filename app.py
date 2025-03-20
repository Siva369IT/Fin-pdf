import streamlit as st
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import fitz  # PyMuPDF for compression

st.set_page_config(page_title="PDF & File Converter", layout="wide")

# Load Logo
st.image("assets/logo1.png", width=120)
st.markdown("<h1 style='text-align:center;'>📄 PDF & File Converter 💚</h1>", unsafe_allow_html=True)

if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = []

if "operation_selected" not in st.session_state:
    st.session_state.operation_selected = None

# Friendly file-type instructions based on operation
def file_instruction(op):
    instructions = {
        "Generate Empty PDF": "No upload required. Just enter number of pages.",
        "Convert Any File to PDF": "Upload PNG, JPG, JPEG, TXT, DOCX, PPTX only.",
        "Images to PDF": "Upload PNG, JPG, JPEG images only.",
        "Extract Pages from PDF": "Upload a single PDF file.",
        "Merge PDFs": "Upload multiple PDF files (at least 2).",
        "Split PDF": "Upload a single PDF file (at least 2 pages).",
        "Compress PDF": "Upload a single PDF file.",
        "Insert Page Numbers": "Upload a single PDF file.",
    }
    return instructions.get(op, "")

st.markdown("<hr>", unsafe_allow_html=True)

# Operation Selection
operations = [
    "Select an operation",
    "Generate Empty PDF",
    "Convert Any File to PDF",
    "Images to PDF",
    "Extract Pages from PDF",
    "Merge PDFs",
    "Split PDF",
    "Compress PDF",
    "Insert Page Numbers"
]

if st.session_state.uploaded_files:
    st.info("**You have uploaded files. Please remove them before changing operation.**")
else:
    operation = st.selectbox("Choose what you'd like to do:", operations, key="operation")
    if operation != "Select an operation":
        st.session_state.operation_selected = operation

if st.session_state.operation_selected:
    st.markdown(f"### 👉 Allowed files: {file_instruction(st.session_state.operation_selected)}")

# Remove Files button
if st.session_state.uploaded_files:
    if st.button("Remove Uploaded Files ❌"):
        st.session_state.uploaded_files = []
        st.rerun()

# Upload Section (only if operation selected)
if st.session_state.operation_selected and st.session_state.operation_selected not in ["Generate Empty PDF"]:
    file_types_map = {
        "Convert Any File to PDF": ["pdf", "png", "jpg", "jpeg", "txt", "docx", "pptx"],
        "Images to PDF": ["png", "jpg", "jpeg"],
        "Extract Pages from PDF": ["pdf"],
        "Merge PDFs": ["pdf"],
        "Split PDF": ["pdf"],
        "Compress PDF": ["pdf"],
        "Insert Page Numbers": ["pdf"]
    }
    allowed_types = file_types_map.get(st.session_state.operation_selected, [])
    uploaded = st.file_uploader("Upload file(s):", type=allowed_types, accept_multiple_files=True)
    if uploaded:
        st.session_state.uploaded_files = uploaded

# Feature Implementations

# 1️⃣ Generate Empty PDF
if st.session_state.operation_selected == "Generate Empty PDF":
    st.subheader("🖨️ Generate Empty PDF")
    pages = st.number_input("Number of pages:maxx-100k", min_value=1, max_value=100000, step=1, value=1)
    if st.button("Generate PDF"):
        output_pdf = BytesIO()
        c = canvas.Canvas(output_pdf, pagesize=letter)
        for i in range(pages):
            c.drawString(250, 750, f"Page {i+1}")
            c.showPage()
        c.save()
        output_pdf.seek(0)
        st.success(f"Generated empty PDF with {pages} pages!")
        st.download_button("Download PDF", data=output_pdf, file_name="empty.pdf", mime="application/pdf")

# 2️⃣ Convert Any File to PDF
elif st.session_state.operation_selected == "Convert Any File to PDF" and st.session_state.uploaded_files:
    st.subheader("♻️ Convert to PDF")
    for file in st.session_state.uploaded_files:
        file_name = file.name.rsplit(".", 1)[0]
        ext = file.name.split(".")[-1].lower()
        pdf_bytes = BytesIO()
        if ext in ["png", "jpg", "jpeg"]:
            img = Image.open(file).convert("RGB")
            img.save(pdf_bytes, format="PDF")
        elif ext == "txt":
            c = canvas.Canvas(pdf_bytes, pagesize=letter)
            text = file.read().decode()
            for line in text.split("\n"):
                c.drawString(100, 750, line)
                c.showPage()
            c.save()
        elif ext == "docx":
            doc = Document(file)
            c = canvas.Canvas(pdf_bytes, pagesize=letter)
            for para in doc.paragraphs:
                c.drawString(100, 750, para.text)
                c.showPage()
            c.save()
        elif ext == "pptx":
            ppt = Presentation(file)
            c = canvas.Canvas(pdf_bytes, pagesize=letter)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        c.drawString(100, 750, shape.text)
                        c.showPage()
            c.save()
        else:
            st.warning(f"Oops! {ext.upper()} not supported for conversion.")
            continue
        pdf_bytes.seek(0)
        st.download_button(f"Download {file_name}.pdf", pdf_bytes, file_name=f"{file_name}.pdf", mime="application/pdf")

# 3️⃣ Images to PDF
elif st.session_state.operation_selected == "Images to PDF" and st.session_state.uploaded_files:
    st.subheader("🏞️ Convert Images to PDF")
    img_files = [Image.open(file).convert("RGB") for file in st.session_state.uploaded_files]
    output_pdf = BytesIO()
    img_files[0].save(output_pdf, save_all=True, append_images=img_files[1:], format="PDF")
    output_pdf.seek(0)
    st.download_button("Download PDF", data=output_pdf, file_name="images_to_pdf.pdf", mime="application/pdf")

# 4️⃣ Extract Pages
elif st.session_state.operation_selected == "Extract Pages from PDF" and st.session_state.uploaded_files:
    st.subheader("🪓 Extract Pages from PDF")
    file = st.session_state.uploaded_files[0]
    if file.type != "application/pdf":
        st.error("Oops! Please upload a valid PDF file for this feature.")
    else:
        pdf = PdfReader(file)
        pages = st.text_input("Enter page numbers (comma-separated):")
        if st.button("Extract Pages"):
            writer = PdfWriter()
            try:
                for p in pages.split(","):
                    idx = int(p.strip()) - 1
                    writer.add_page(pdf.pages[idx])
                output = BytesIO()
                writer.write(output)
                output.seek(0)
                st.download_button("Download Extracted PDF", data=output, file_name="extracted_pages.pdf", mime="application/pdf")
            except:
                st.error("Invalid page numbers provided.")

# 5️⃣ Merge PDFs
elif st.session_state.operation_selected == "Merge PDFs" and st.session_state.uploaded_files:
    st.subheader("📄+📄 Merge PDFs")
    if len(st.session_state.uploaded_files) < 2:
        st.warning("Upload at least two PDFs to merge.")
    else:
        writer = PdfWriter()
        for file in st.session_state.uploaded_files:
            reader = PdfReader(file)
            for page in reader.pages:
                writer.add_page(page)
        output = BytesIO()
        writer.write(output)
        output.seek(0)
        st.download_button("Download Merged PDF", data=output, file_name="merged.pdf", mime="application/pdf")

# 6️⃣ Split PDF
elif st.session_state.operation_selected == "Split PDF" and st.session_state.uploaded_files:
    st.subheader("📑 Split PDF into 2 parts")
    file = st.session_state.uploaded_files[0]
    pdf = PdfReader(file)
    if len(pdf.pages) <= 1:
        st.warning("Cannot split a single-page PDF.")
    else:
        split_at = st.number_input("Enter split page number:", min_value=1, max_value=len(pdf.pages)-1, step=1)
        if st.button("Split"):
            part1, part2 = PdfWriter(), PdfWriter()
            for i in range(len(pdf.pages)):
                if i < split_at:
                    part1.add_page(pdf.pages[i])
                else:
                    part2.add_page(pdf.pages[i])
            out1, out2 = BytesIO(), BytesIO()
            part1.write(out1)
            part2.write(out2)
            out1.seek(0)
            out2.seek(0)
            st.download_button("Download Part 1", data=out1, file_name="split_part1.pdf", mime="application/pdf")
            st.download_button("Download Part 2", data=out2, file_name="split_part2.pdf", mime="application/pdf")

# 7️⃣ Compress PDF
elif st.session_state.operation_selected == "Compress PDF" and st.session_state.uploaded_files:
    st.subheader("📉 Compress PDF")
    file = st.session_state.uploaded_files[0]
    pdf = fitz.open(stream=file.getvalue(), filetype="pdf")
    output = BytesIO()
    pdf.save(output, garbage=4, deflate=True)
    output.seek(0)
    st.download_button("Download Compressed PDF", data=output, file_name="compressed.pdf", mime="application/pdf")

# 8️⃣ Insert Page Numbers
elif st.session_state.operation_selected == "Insert Page Numbers" and st.session_state.uploaded_files:
    st.subheader("📝 Insert Page Numbers")
    file = st.session_state.uploaded_files[0]
    pdf = PdfReader(file)
    writer = PdfWriter()
    for i, page in enumerate(pdf.pages):
        packet = BytesIO()
        c = canvas.Canvas(packet, pagesize=letter)
        c.drawString(500, 20, f"Page {i + 1}")
        c.save()
        packet.seek(0)
        overlay = PdfReader(packet)
        page.merge_page(overlay.pages[0])
        writer.add_page(page)
    output = BytesIO()
    writer.write(output)
    output.seek(0)
    st.download_button("Download PDF with Page Numbers", data=output, file_name="numbered.pdf", mime="application/pdf")

# Footer
st.markdown('<hr>', unsafe_allow_html=True)
st.markdown(
    '<div style="text-align:center; font-size:small;">'
    '© Pavan Sri Sai Mondem | Siva Satyamsetti | Uma Satya Mounika Sapireddy | '
    'Bhuvaneswari Devi Seru | Chandu Meela | Trainees from techwing 🧡 </div>', 
    unsafe_allow_html=True
                                           )
